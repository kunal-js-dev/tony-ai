"""PowerPoint slide extraction using python-pptx."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_ppt(filepath: str | Path) -> dict[str, Any]:
    """Extract text from all slides in a PowerPoint file."""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"Presentation not found: {path}")

    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    text_parts: list[str] = []

    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(row_cells))

        slide_content = "\n".join(slide_texts)
        slides.append(
            {
                "slide_number": i + 1,
                "text": slide_content,
                "shape_count": len(slide.shapes),
            }
        )
        if slide_content:
            text_parts.append(f"--- Slide {i + 1} ---\n{slide_content}")

    full_text = "\n\n".join(text_parts)

    return {
        "type": "ppt",
        "title": path.stem,
        "slide_count": len(slides),
        "word_count": len(full_text.split()) if full_text else 0,
        "text": full_text,
        "slides": slides,
        "summary": {
            "slides": len(slides),
            "words": len(full_text.split()) if full_text else 0,
        },
    }
